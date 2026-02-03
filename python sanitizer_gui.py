# python sanitizer_gui.py
# pip install pyinstallerでexe化して使用することを想定。実行するとアプリが起動し、ディレクトリパスとサニタイズ対象フレーズの設定を求められる。フレーズは10個まで登録可能。実行ボタンを押下するとディレクトリ配下のファイルに対してサニタイズ処理が実行される。

import os
import re
import threading
import queue
from dataclasses import dataclass
from pathlib import Path
import tkinter as tk
from tkinter import ttk, filedialog, messagebox

# Office
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook


REPLACEMENT = "xxx"

# 対象拡張子（必要なら追加してください）
SUPPORTED_EXTS = {
    ".docx",             # Word
    ".pptx",             # PowerPoint
    ".xlsx", ".xlsm",    # Excel
    ".txt",              # メモ帳
    ".md", ".markdown",  # Markdown
    ".html", ".htm"      # HTML
}

# 既定の自動検出パターン（必要に応じて強化/緩和してください）
DEFAULT_PATTERNS = [
    # IPv4（厳密ではないが実用的）
    r"\b(?:(?:25[0-5]|2[0-4]\d|1?\d?\d)\.){3}(?:25[0-5]|2[0-4]\d|1?\d?\d)\b",

    # IPv6（簡易）
    r"\b(?:[0-9A-Fa-f]{1,4}:){2,7}[0-9A-Fa-f]{1,4}\b",

    # Windows パス（例: C:\Users\name\file.txt）※スペースや日本語をある程度許容
    r"\b[A-Za-z]:\\(?:[^\\/:*?\"<>|\r\n]+\\)*[^\\/:*?\"<>|\r\n]+\b",

    # UNC パス（例: \\server\share\dir\file）
    r"\\\\[^\s\\/:*?\"<>|\r\n]+\\[^\s\\/:*?\"<>|\r\n]+(?:\\[^\s\\/:*?\"<>|\r\n]+)*",

    # Unix系パス（例: /home/user/file）※誤検出が嫌なら削除
    r"(?<!\w)/(?:[\w.\-]+/)+[\w.\-]+",

    # URL
    r"\bhttps?://[^\s<>\"]+\b",

    # Email
    r"\b[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}\b",

    # MACアドレス
    r"\b(?:[0-9A-Fa-f]{2}[:-]){5}[0-9A-Fa-f]{2}\b",
]


@dataclass
class SanitizeConfig:
    base_dir: Path
    out_dir: Path
    overwrite: bool
    ignore_case: bool
    user_phrases: list[str]


class Sanitizer:
    def __init__(self, config: SanitizeConfig):
        self.config = config

        flags = re.IGNORECASE if config.ignore_case else 0

        # ユーザー指定フレーズ（リテラル一致を安全にするため re.escape）
        self.user_regex = None
        cleaned = [p for p in (s.strip() for s in config.user_phrases) if p]
        if cleaned:
            # 長いものから先にマッチさせたいのでソート（部分一致の置換順事故を減らす）
            cleaned.sort(key=len, reverse=True)
            self.user_regex = re.compile("|".join(re.escape(p) for p in cleaned), flags=flags)

        # 既定の自動検出
        self.default_res = [re.compile(p, flags=flags) for p in DEFAULT_PATTERNS]

    def sanitize_text(self, text: str) -> tuple[str, int]:
        """
        text をサニタイズし、(サニタイズ後テキスト, 置換回数合計) を返す
        """
        if not text:
            return text, 0

        total = 0
        out = text

        # ユーザー指定フレーズ
        if self.user_regex is not None:
            out, n = self.user_regex.subn(REPLACEMENT, out)
            total += n

        # 既定パターン
        for rx in self.default_res:
            out, n = rx.subn(REPLACEMENT, out)
            total += n

        return out, total

    # ---------- ファイル別処理 ----------

    def handle_text_file(self, src: Path, dst: Path) -> tuple[int, int]:
        """
        (変更が入ったファイル=1/0, 置換回数)
        """
        # 文字コードはまずutf-8、ダメならcp932（Windowsのメモ帳想定）
        raw = None
        for enc in ("utf-8", "cp932", "utf-8-sig"):
            try:
                raw = src.read_text(encoding=enc)
                break
            except UnicodeDecodeError:
                continue

        if raw is None:
            # 最後の手段: バイナリとして読んで置換しない（誤破壊を避ける）
            return 0, 0

        sanitized, n = self.sanitize_text(raw)
        changed = 1 if sanitized != raw else 0

        # 出力
        dst.parent.mkdir(parents=True, exist_ok=True)
        dst.write_text(sanitized, encoding="utf-8")
        return changed, n

    def handle_docx(self, src: Path, dst: Path) -> tuple[int, int]:
        """
        Word(docx) の本文・表・ヘッダー/フッターのテキストを置換
        ※テキストボックス等は python-docx の制約で拾えない場合があります
        """
        doc = Document(str(src))
        total = 0
        changed_flag = 0

        def sanitize_paragraph(par):
            nonlocal total, changed_flag
            for run in par.runs:
                new_text, n = self.sanitize_text(run.text)
                if new_text != run.text:
                    run.text = new_text
                    changed_flag = 1
                total += n

        def sanitize_table(table):
            for row in table.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        sanitize_paragraph(p)
                    for t in cell.tables:
                        sanitize_table(t)

        # 本文
        for p in doc.paragraphs:
            sanitize_paragraph(p)
        for t in doc.tables:
            sanitize_table(t)

        # ヘッダー/フッター
        for section in doc.sections:
            header = section.header
            footer = section.footer
            for p in header.paragraphs:
                sanitize_paragraph(p)
            for t in header.tables:
                sanitize_table(t)
            for p in footer.paragraphs:
                sanitize_paragraph(p)
            for t in footer.tables:
                sanitize_table(t)

        dst.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(dst))
        return changed_flag, total

    def handle_pptx(self, src: Path, dst: Path) -> tuple[int, int]:
        """
        PowerPoint(pptx) のスライド内テキストとノートを置換
        """
        prs = Presentation(str(src))
        total = 0
        changed_flag = 0

        def sanitize_textframe(tf):
            nonlocal total, changed_flag
            for p in tf.paragraphs:
                for run in p.runs:
                    new_text, n = self.sanitize_text(run.text)
                    if new_text != run.text:
                        run.text = new_text
                        changed_flag = 1
                    total += n

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    sanitize_textframe(shape.text_frame)

            # ノート
            if slide.has_notes_slide and slide.notes_slide is not None:
                ns = slide.notes_slide
                for shape in ns.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        sanitize_textframe(shape.text_frame)

        dst.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(dst))
        return changed_flag, total

    def handle_xlsx(self, src: Path, dst: Path) -> tuple[int, int]:
        """
        Excel(xlsx/xlsm) のセル文字列・コメントを置換
        """
        wb = load_workbook(str(src), keep_vba=src.suffix.lower() == ".xlsm")
        total = 0
        changed_flag = 0

        for ws in wb.worksheets:
            # セル値
            for row in ws.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str):
                        new_text, n = self.sanitize_text(cell.value)
                        if new_text != cell.value:
                            cell.value = new_text
                            changed_flag = 1
                        total += n

                    # コメント
                    if cell.comment is not None and isinstance(cell.comment.text, str):
                        new_text, n = self.sanitize_text(cell.comment.text)
                        if new_text != cell.comment.text:
                            cell.comment.text = new_text
                            changed_flag = 1
                        total += n

        dst.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(dst))
        return changed_flag, total

    # ---------- 走査 ----------

    def sanitize_directory(self, progress_cb=None, log_cb=None) -> dict:
        """
        ディレクトリ配下を再帰的に処理し、結果統計を返す
        """
        base = self.config.base_dir.resolve()

        # 出力ディレクトリが入力配下にあると無限ループになるので回避
        out_dir = self.config.out_dir.resolve()
        if not self.config.overwrite and str(out_dir).startswith(str(base) + os.sep):
            raise ValueError("出力先が入力ディレクトリの配下にあります。別の出力先にしてください。")

        files = []
        for root, dirs, filenames in os.walk(base):
            root_path = Path(root)

            # 出力フォルダは走査しない（上書きしない場合）
            if not self.config.overwrite:
                # 走査中に out_dir が見つかったらスキップ
                dirs[:] = [d for d in dirs if (root_path / d).resolve() != out_dir]

            for fn in filenames:
                p = root_path / fn
                if p.suffix.lower() in SUPPORTED_EXTS:
                    files.append(p)

        total_files = len(files)
        stat = {
            "scanned": total_files,
            "changed_files": 0,
            "total_replacements": 0,
            "skipped": 0,
            "errors": 0,
        }

        def log(msg: str):
            if log_cb:
                log_cb(msg)

        for i, src in enumerate(files, start=1):
            rel = src.relative_to(base)

            if self.config.overwrite:
                dst = src
            else:
                dst = out_dir / rel

            try:
                ext = src.suffix.lower()
                if ext in (".txt", ".md", ".markdown", ".html", ".htm"):
                    changed, reps = self.handle_text_file(src, dst)
                elif ext == ".docx":
                    changed, reps = self.handle_docx(src, dst)
                elif ext == ".pptx":
                    changed, reps = self.handle_pptx(src, dst)
                elif ext in (".xlsx", ".xlsm"):
                    changed, reps = self.handle_xlsx(src, dst)
                else:
                    stat["skipped"] += 1
                    continue

                stat["changed_files"] += changed
                stat["total_replacements"] += reps

                log(f"[{i}/{total_files}] OK: {rel}  (置換 {reps} 回)")
            except Exception as e:
                stat["errors"] += 1
                log(f"[{i}/{total_files}] ERROR: {rel}  ({e})")

            if progress_cb:
                progress_cb(i, total_files)

        return stat


class App(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("ファイル自動サニタイズ（Word/PPT/Excel/txt/md/html）")
        self.geometry("900x650")

        self.q = queue.Queue()
        self.worker_thread = None

        # 入力変数
        self.dir_var = tk.StringVar()
        self.overwrite_var = tk.BooleanVar(value=False)
        self.ignore_case_var = tk.BooleanVar(value=True)

        # UI
        self._build_ui()
        self._poll_queue()

    def _build_ui(self):
        pad = {"padx": 10, "pady": 6}

        frm_top = ttk.LabelFrame(self, text="対象ディレクトリ")
        frm_top.pack(fill="x", **pad)

        row = ttk.Frame(frm_top)
        row.pack(fill="x", padx=10, pady=8)

        ttk.Label(row, text="ディレクトリ: ").pack(side="left")
        entry = ttk.Entry(row, textvariable=self.dir_var)
        entry.pack(side="left", fill="x", expand=True, padx=6)

        ttk.Button(row, text="参照...", command=self.choose_dir).pack(side="left")

        opt = ttk.Frame(frm_top)
        opt.pack(fill="x", padx=10, pady=6)

        ttk.Checkbutton(opt, text="上書きする（注意）", variable=self.overwrite_var).pack(side="left")
        ttk.Checkbutton(opt, text="大文字小文字を無視（入力文章・既定パターン）", variable=self.ignore_case_var).pack(side="left", padx=18)

        frm_phrases = ttk.LabelFrame(self, text="置換したい文章（最大10個・空欄OK）")
        frm_phrases.pack(fill="x", **pad)

        self.phrase_vars = []
        grid = ttk.Frame(frm_phrases)
        grid.pack(fill="x", padx=10, pady=8)

        for i in range(10):
            v = tk.StringVar()
            self.phrase_vars.append(v)
            r = i // 2
            c = i % 2
            cell = ttk.Frame(grid)
            cell.grid(row=r, column=c, sticky="ew", padx=6, pady=4)
            ttk.Label(cell, text=f"{i+1}:").pack(side="left")
            ttk.Entry(cell, textvariable=v, width=45).pack(side="left", fill="x", expand=True, padx=6)

        grid.columnconfigure(0, weight=1)
        grid.columnconfigure(1, weight=1)

        frm_run = ttk.Frame(self)
        frm_run.pack(fill="x", **pad)

        self.btn_run = ttk.Button(frm_run, text="サニタイズ実行", command=self.on_run)
        self.btn_run.pack(side="left")

        self.progress = ttk.Progressbar(frm_run, orient="horizontal", mode="determinate")
        self.progress.pack(side="left", fill="x", expand=True, padx=10)

        self.lbl_status = ttk.Label(frm_run, text="待機中")
        self.lbl_status.pack(side="left")

        frm_log = ttk.LabelFrame(self, text="ログ")
        frm_log.pack(fill="both", expand=True, **pad)

        self.txt_log = tk.Text(frm_log, height=18)
        self.txt_log.pack(fill="both", expand=True, padx=10, pady=8)

        note = (
            "既定の自動検出: IP(IPv4/IPv6), Windowsパス, UNCパス, URL, Email, MAC\n"
            f"置換文字列: {REPLACEMENT}\n"
            "※ Officeはテキスト部分のみ。図形/埋め込み/一部テキストボックス等は拾えない場合があります。\n"
        )
        self.txt_log.insert("end", note + "\n")
        self.txt_log.configure(state="disabled")

    def choose_dir(self):
        d = filedialog.askdirectory()
        if d:
            self.dir_var.set(d)

    def log(self, msg: str):
        self.txt_log.configure(state="normal")
        self.txt_log.insert("end", msg + "\n")
        self.txt_log.see("end")
        self.txt_log.configure(state="disabled")

    def set_status(self, msg: str):
        self.lbl_status.configure(text=msg)

    def on_run(self):
        base_dir = self.dir_var.get().strip()
        if not base_dir:
            messagebox.showerror("エラー", "ディレクトリを指定してください。")
            return

        base = Path(base_dir)
        if not base.exists() or not base.is_dir():
            messagebox.showerror("エラー", "正しいディレクトリを指定してください。")
            return

        overwrite = bool(self.overwrite_var.get())
        ignore_case = bool(self.ignore_case_var.get())
        phrases = [v.get() for v in self.phrase_vars]

        out_dir = base if overwrite else Path(str(base) + "_sanitized")

        cfg = SanitizeConfig(
            base_dir=base,
            out_dir=out_dir,
            overwrite=overwrite,
            ignore_case=ignore_case,
            user_phrases=phrases
        )

        self.btn_run.configure(state="disabled")
        self.progress["value"] = 0
        self.set_status("実行中...")

        # ワーカースレッドで実行
        self.worker_thread = threading.Thread(target=self._worker, args=(cfg,), daemon=True)
        self.worker_thread.start()

    def _worker(self, cfg: SanitizeConfig):
        try:
            sanitizer = Sanitizer(cfg)

            def progress_cb(done, total):
                self.q.put(("progress", done, total))

            def log_cb(msg):
                self.q.put(("log", msg))

            stat = sanitizer.sanitize_directory(progress_cb=progress_cb, log_cb=log_cb)
            self.q.put(("done", stat))
        except Exception as e:
            self.q.put(("fatal", str(e)))

    def _poll_queue(self):
        try:
            while True:
                item = self.q.get_nowait()
                kind = item[0]

                if kind == "log":
                    self.log(item[1])

                elif kind == "progress":
                    done, total = item[1], item[2]
                    self.progress["maximum"] = max(total, 1)
                    self.progress["value"] = done

                elif kind == "done":
                    stat = item[1]
                    self.set_status("完了")
                    self.btn_run.configure(state="normal")

                    msg = (
                        f"完了しました。\n"
                        f"対象ファイル数: {stat['scanned']}\n"
                        f"変更が入ったファイル: {stat['changed_files']}\n"
                        f"置換回数合計: {stat['total_replacements']}\n"
                        f"スキップ: {stat['skipped']}\n"
                        f"エラー: {stat['errors']}\n"
                    )
                    self.log("\n" + msg)
                    if not cfg.overwrite:
                        self.log(f"出力先: {cfg.out_dir}\n")

                    messagebox.showinfo("完了", msg)

                elif kind == "fatal":
                    self.set_status("エラー")
                    self.btn_run.configure(state="normal")
                    messagebox.showerror("エラー", item[1])

        except queue.Empty:
            pass

        self.after(100, self._poll_queue)


def main():
    app = App()
    app.mainloop()


if __name__ == "__main__":
    main()
